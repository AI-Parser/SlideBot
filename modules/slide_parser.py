from pptx import Presentation
import google.generativeai as genai

def parse_pptx(file_path, api_key):
    prs = Presentation(file_path)
    slides_data = []

    for slide in prs.slides:
        title = ''
        content = []
        if slide.shapes:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape == slide.shapes[0]:
                        title = shape.text
                    else:
                        content.append(shape.text)
        slides_data.append({
            "title": title,
            "content": "\n".join(content)
        })

    return generate_pitch_notes(slides_data, api_key)

def generate_pitch_notes(slides_data, api_key):
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("models/gemini-1.5-flash")

    pitch_notes = []
    for slide in slides_data:
        prompt = f"""You are helping prepare speaker notes for a pitch.
Here is the slide content:

Title: {slide['title']}
Content:
{slide['content']}

Summarize this into concise, persuasive pitch notes for a presenter to speak during a presentation.
Make it sound confident and clear, like it's being spoken aloud.
Avoid repeating the title unless necessary."""

        response = model.generate_content(prompt)
        pitch_notes.append({
            "title": slide['title'],
            "original": slide['content'],
            "content": response.text.strip()
        })

    print(pitch_notes)

    return pitch_notes